from flask import Flask, render_template, request, send_file
from pptx import Presentation
import random

app = Flask(__name__)

def generate_dummy_slides(topic):
    """Generate simple 5-slide content without AI/API"""
    slides = []
    for i in range(1, 6):
        title = f"Slide {i}: {topic}"
        points = [
            f"Introduction to {topic} - point {i}",
            f"Key aspect of {topic} - {random.choice(['overview', 'importance', 'benefit'])}",
            f"Further details about {topic} - point {i}"
        ]
        slides.append((title, points))
    return slides

def create_ppt(slides, filename="presentation.pptx"):
    prs = Presentation()
    for title, points in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        for p in points:
            body.add_paragraph(p)
    prs.save(filename)
    return filename

@app.route("/")
def index():
    return render_template("index.html")

@app.route("/generate", methods=["POST"])
def generate():
    topic = request.form["topic"]
    slides = generate_dummy_slides(topic)
    filename = "output.pptx"
    create_ppt(slides, filename)
    return send_file(filename, as_attachment=True)

if __name__ == "__main__":
    app.run(debug=True)
